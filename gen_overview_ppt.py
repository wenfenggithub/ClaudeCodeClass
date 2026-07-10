#!/usr/bin/env python3
"""生成一页课程总览 PPT（Apple 风格），含 GitHub 网址。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ---- 调色板 ----
BG      = RGBColor(0xFB, 0xFB, 0xFD)
FG      = RGBColor(0x1D, 0x1D, 0x1F)
FG_DIM  = RGBColor(0x6E, 0x6E, 0x73)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
BORDER  = RGBColor(0xE3, 0xE3, 0xE8)

COLORS = [
    RGBColor(0xFF, 0x3B, 0x30),  # 1 red
    RGBColor(0xFF, 0x95, 0x00),  # 2 orange
    RGBColor(0x28, 0xA7, 0x45),  # 3 green
    RGBColor(0x00, 0x71, 0xE3),  # 4 blue
    RGBColor(0xAF, 0x52, 0xDE),  # 5 purple
    RGBColor(0x5A, 0xC8, 0xFA),  # 6 cyan
    RGBColor(0xFF, 0x2D, 0x55),  # 7 pink
]
ACCENT  = RGBColor(0x00, 0x71, 0xE3)

FONT = "PingFang SC"

def no_line(shape):
    shape.line.fill.background()

def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_radius(shape, frac=0.12):
    try:
        shape.adjustments[0] = frac
    except Exception:
        pass

def add_text(left, top, width, height, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0):
    """runs: list of (text, size, bold, color, [font])"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    if space_after:
        p.space_after = Pt(space_after)
    for r in runs:
        text, size, bold, color = r[0], r[1], r[2], r[3]
        fname = r[4] if len(r) > 4 else FONT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = fname
    return tb

# ---- 背景 ----
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
solid(bg, BG); no_line(bg)
bg.shadow.inherit = False

# ---- 顶部标题区 ----
MARGIN = Inches(0.62)
add_text(MARGIN, Inches(0.42), Inches(12), Inches(0.4),
         [("CLAUDE CODE · 零基础 AI 编程课", 12, True, FG_DIM)])

add_text(MARGIN, Inches(0.74), Inches(12.2), Inches(0.9),
         [("AI 分享", 34, True, FG)])

add_text(MARGIN, Inches(1.52), Inches(11.6), Inches(0.7),
         [("面向零基础成年人的 AI 编程实战课 —— 你需要的不是写代码，而是学会指挥 AI。",
           14, False, FG_DIM)], line_spacing=1.2)

# ---- 卡片网格 4 x 2 ----
CARDS = [
    ("1", "\"哇\"时刻 + 环境搭建", "见证 AI 把想法变网页，装好开发环境。", "点燃 · 起步"),
    ("2", "AI 思维模型 + 初次动手", "任务书 → 生成 → 评估 → 优化的核心循环。", "核心循环"),
    ("3", "Harness Engineering", "给 AI 装\"缰绳\"：用引导与反馈保稳定。", "引导 · 反馈"),
    ("4", "客户端、服务器与联网应用", "从纯 HTML 走向真正的联网 Web 应用。", "Web 应用"),
    ("5", "Claude Code Skill", "为 AI 打造可复用的\"技能\"，做得更专业。", "进阶能力"),
    ("6", "AI 数据分析与可视化", "从数据里找答案，把答案用图表讲清楚。", "数据 · 图表"),
    ("7", "迁移到 Supabase 数据库", "从 JSON 文件升级到云数据库，可上线。", "数据库 · 上线"),
]

cols, rows = 4, 2
gx, gy = Inches(0.28), Inches(0.28)
grid_left = MARGIN
grid_top = Inches(2.42)
grid_w = SW - 2 * MARGIN
card_w = (grid_w - gx * (cols - 1)) / cols
card_h = Inches(1.92)

for i, (num, title, desc, tag) in enumerate(CARDS):
    r, c = divmod(i, cols)
    x = grid_left + c * (card_w + gx)
    y = grid_top + r * (card_h + gy)
    color = COLORS[i]

    # 卡片底
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    solid(card, CARD_BG); set_radius(card, 0.09)
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.75)
    card.shadow.inherit = False

    # 顶部彩条
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.14),
                                 y + Inches(0.12), Inches(0.42), Inches(0.42))
    solid(bar, color); no_line(bar); set_radius(bar, 0.28)
    bar.shadow.inherit = False
    btf = bar.text_frame
    btf.margin_left = 0; btf.margin_right = 0
    btf.margin_top = 0; btf.margin_bottom = 0
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = num
    br.font.size = Pt(18); br.font.bold = True
    br.font.color.rgb = WHITE; br.font.name = FONT

    pad = Inches(0.2)
    add_text(x + pad, y + Inches(0.66), card_w - 2 * pad, Inches(0.5),
             [(title, 14.5, True, FG)], line_spacing=1.05)
    add_text(x + pad, y + Inches(1.14), card_w - 2 * pad, Inches(0.55),
             [(desc, 10.5, False, FG_DIM)], line_spacing=1.15)
    # tag
    add_text(x + pad, y + card_h - Inches(0.33), card_w - 2 * pad, Inches(0.25),
             [(tag, 9.5, True, color)])

# ---- 第 8 格：GitHub 链接卡 ----
i = 7
r, c = divmod(i, cols)
x = grid_left + c * (card_w + gx)
y = grid_top + r * (card_h + gy)
gh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
solid(gh, FG); set_radius(gh, 0.09); no_line(gh)
gh.shadow.inherit = False
pad = Inches(0.2)
add_text(x + pad, y + Inches(0.22), card_w - 2 * pad, Inches(0.4),
         [("</>  课程代码仓库", 13, True, WHITE)])
add_text(x + pad, y + Inches(0.74), card_w - 2 * pad, Inches(0.4),
         [("全部课件 · 作业 · 示例", 10.5, False, RGBColor(0xC7,0xC7,0xCC))])
add_text(x + pad, y + Inches(1.18), card_w - 2 * pad, Inches(0.65),
         [("github.com/wenfenggithub/", 10.5, True, RGBColor(0x5A,0xC8,0xFA)),
          ("ClaudeCodeClass", 10.5, True, RGBColor(0x5A,0xC8,0xFA))], line_spacing=1.1)

# 让 GitHub 卡文字可点击（超链接放到不可见的形状上）
link = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
link.fill.background(); no_line(link); link.shadow.inherit = False
link.click_action.hyperlink.address = "https://github.com/wenfenggithub/ClaudeCodeClass"

# ---- 底部信息条 ----
add_text(MARGIN, SH - Inches(0.5), Inches(12), Inches(0.35),
         [("7 次课 · 每次 2 小时", 11, True, FG),
          ("      主线：从想法到上线的完整闭环", 11, False, FG_DIM)])

out = "/Users/wenfengzhu/ClaudeProjects/ClaudeClass/course-overview.pptx"
prs.save(out)
print("saved:", out)
